import json
import os
import re
import sys
import time

from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

OUTPUT_DIR = "output"
MODEL = "gemini-3.1-flash-lite"
SLEEP_BETWEEN_TOPICS = 5
THAI_FONT = "Leelawadee UI"


def set_thai_font(run, size_pt=None, bold=None, color=None):
    """Set font on a run so Thai text renders correctly (latin + complex-script)."""
    from pptx.oxml.ns import qn
    run.font.name = THAI_FONT
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", THAI_FONT)


def set_thai_font_frame(text_frame, size_pt=None, bold=None, color=None):
    for p in text_frame.paragraphs:
        for run in p.runs:
            set_thai_font(run, size_pt=size_pt, bold=bold, color=color)

API_KEY = os.environ.get("GEMINI_API_KEY")
if not API_KEY:
    print("ERROR: GEMINI_API_KEY environment variable not set.")
    sys.exit(1)

client = genai.Client(api_key=API_KEY)

SLIDE_PROMPT_TEMPLATE = """คุณเป็นผู้เชี่ยวชาญด้านการสอน CCNP ENCOR 350-401 กำลังออกแบบสไลด์สอนหัวข้อ "{topic}"
จากเนื้อหาสรุปด้านล่างนี้ (ภาษาไทย):

---
{summary}
---

จงแปลงเนื้อหานี้เป็นโครงสไลด์สำหรับสอน (ไม่ใช่แค่ก๊อปข้อความมาแปะ) ภาษาไทย โดยออกแบบลำดับการสอนให้เข้าใจง่าย:
1. Slide เปิดหัวข้อ (Title + เกริ่นว่าเรียนเรื่องอะไร ทำไมสำคัญ)
2. Slide อธิบาย concept หลัก
3. Slide diagram แสดงขั้นตอน/การทำงาน/topology (ถ้าเนื้อหาเหมาะกับการวาดเป็นภาพ เช่น process flow, state machine, network topology, packet flow) — ใส่ "diagram" ด้วย
4. Slide อธิบายการทำงาน (How it works) อาจมีหลาย slide ถ้าเนื้อหายาว
5. Slide ตัวอย่าง config จริง (ถ้ามีในเนื้อหา)
6. Slide เปรียบเทียบ/common pitfalls (ถ้ามีในเนื้อหา) — ถ้าเป็นการเปรียบเทียบ 2 สิ่ง อาจใส่ "diagram" แบบ compare ด้วย
7. Slide สรุป Key points สำหรับสอบ

ตอบเป็น JSON เท่านั้น (ไม่ต้องมีคำอธิบายอื่นนอก JSON) ในรูปแบบ:
{{
  "slides": [
    {{
      "title": "หัวข้อสไลด์ (สั้น กระชับ)",
      "bullets": ["bullet 1", "bullet 2", "..."],
      "notes": "คำอธิบายเพิ่มเติมสำหรับคนสอน (speaker notes) อธิบายละเอียดกว่าบน slide",
      "diagram": {{
        "nodes": ["กล่องที่ 1", "กล่องที่ 2", "กล่องที่ 3"],
        "edges": [[0, 1], [1, 2]]
      }}
    }}
  ]
}}

ข้อกำหนด:
- bullet แต่ละอันสั้น กระชับ ไม่เกิน 1-2 บรรทัด ไม่ใส่ยาวเป็นพารากราฟ
- slide ละไม่เกิน 5-6 bullets
- ถ้ามี config ให้ใส่ในรูป bullet ที่เป็นคำสั่งจริง (ไม่ต้องห่อ code block)
- รวมทั้งหมดควรมีประมาณ 6-10 slides ต่อหัวข้อ
- "diagram" เป็น field เสริม ใส่เฉพาะ slide ที่เนื้อหาเหมาะกับภาพ flow/topology/comparison เท่านั้น (ไม่ต้องใส่ทุก slide) ถ้า slide ไหนไม่เหมาะให้ใส่ "diagram": null
- "nodes" ควรมีไม่เกิน 6 กล่อง ข้อความในกล่องสั้นมาก (ไม่เกิน 4-5 คำ)
- "edges" คือคู่ index ของ nodes ที่มีลูกศรชี้จาก -> ไป (เช่น [0,1] หมายถึงลูกศรจาก nodes[0] ไป nodes[1])
"""


def slugify_check(folder_name: str) -> str:
    return folder_name


def call_gemini(topic: str, summary: str) -> dict:
    prompt = SLIDE_PROMPT_TEMPLATE.format(topic=topic, summary=summary)
    response = client.models.generate_content(
        model=MODEL,
        contents=prompt,
        config={"response_mime_type": "application/json"},
    )
    text = response.text
    return json.loads(text)


def add_title_slide(prs: Presentation, topic: str):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = topic
    set_thai_font_frame(slide.shapes.title.text_frame, bold=True)
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "CCNP ENCOR 350-401"
        set_thai_font_frame(slide.placeholders[1].text_frame)
    return slide


def add_content_slide(prs: Presentation, title: str, bullets: list, notes: str, diagram: dict = None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    set_thai_font_frame(slide.shapes.title.text_frame, bold=True)

    has_diagram = bool(diagram and diagram.get("nodes"))

    body = slide.placeholders[1]
    if has_diagram:
        # shrink the bullet box to make room for the diagram below it
        body.height = Inches(2.6)

    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(bullet)
        p.level = 0
        for run in p.runs:
            set_thai_font(run, size_pt=(18 if has_diagram else 20))

    if has_diagram:
        add_diagram_to_slide(slide, diagram.get("nodes", []), diagram.get("edges", []), top_in=3.4)

    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_diagram_to_slide(slide, nodes: list, edges: list, top_in: float = 3.6):
    """Draw simple boxes-and-arrows flow diagram onto an existing slide."""
    if not nodes:
        return

    n = len(nodes)
    slide_width_in = 13.333
    margin_in = 0.6
    usable_width = slide_width_in - 2 * margin_in
    box_w = min(2.6, usable_width / n - 0.3)
    box_h = 1.0
    gap = (usable_width - n * box_w) / max(n - 1, 1) if n > 1 else 0

    box_shapes = []
    for i, text in enumerate(nodes):
        left = Inches(margin_in + i * (box_w + gap))
        top = Inches(top_in)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(box_w), Inches(box_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x2E, 0x75, 0xB6)
        shape.line.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.text = str(text)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            set_thai_font(run, size_pt=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        box_shapes.append(shape)

    for edge in edges:
        try:
            a, b = edge[0], edge[1]
            shape_a, shape_b = box_shapes[a], box_shapes[b]
        except (IndexError, TypeError, ValueError):
            continue
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            shape_a.left + shape_a.width, shape_a.top + shape_a.height // 2,
            shape_b.left, shape_b.top + shape_b.height // 2,
        )
        connector.line.color.rgb = RGBColor(0x40, 0x40, 0x40)
        connector.line.width = Pt(2)


def build_pptx(topic: str, slides_data: list, out_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, topic)

    for s in slides_data:
        add_content_slide(
            prs,
            s.get("title", ""),
            s.get("bullets", []),
            s.get("notes", ""),
            s.get("diagram"),
        )

    prs.save(out_path)


def process_topic(folder: str):
    md_path = os.path.join(OUTPUT_DIR, folder, "summary_th.md")
    out_path = os.path.join(OUTPUT_DIR, folder, "slide_gemini.pptx")

    if not os.path.exists(md_path):
        print(f"  SKIP (no summary): {folder}")
        return
    if os.path.exists(out_path):
        print(f"  SKIP (already exists): {folder}")
        return

    with open(md_path, "r", encoding="utf-8") as f:
        raw = f.read()

    lines = raw.split("\n")
    topic = lines[0][2:].strip() if lines[0].startswith("# ") else folder
    summary = "\n".join(lines[1:])

    print(f"\n[{folder}] {topic}")
    print("  Calling Gemini...")
    try:
        data = call_gemini(topic, summary)
        slides_data = data.get("slides", [])
        if not slides_data:
            print("  ERROR: Gemini returned no slides")
            return
        build_pptx(topic, slides_data, out_path)
        print(f"  OK: {len(slides_data) + 1} slides -> {out_path}")
    except Exception as e:
        print(f"  ERROR: {e}")


def main():
    start_id = None
    if "--start-id" in sys.argv:
        start_id = sys.argv[sys.argv.index("--start-id") + 1]

    folders = sorted(
        f for f in os.listdir(OUTPUT_DIR)
        if os.path.isdir(os.path.join(OUTPUT_DIR, f))
    )

    started = start_id is None
    total = len(folders)

    for i, folder in enumerate(folders):
        if not started:
            if folder.startswith(start_id):
                started = True
            else:
                continue

        process_topic(folder)

        remaining = total - i - 1
        if remaining > 0:
            time.sleep(SLEEP_BETWEEN_TOPICS)

    print("\nAll slide decks complete!")


if __name__ == "__main__":
    main()
